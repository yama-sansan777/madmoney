# -*- coding: utf-8 -*-
"""
マッドマネー要約スライド 汎用生成スクリプト
使い方：
    python generate_pptx.py plans/plan_YYYYMMDD_タイトル.md
"""

import sys
import re
import os
sys.path.insert(0, r'C:\Users\PC-YM\AppData\Roaming\Python\Python313\site-packages')

from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ============================================================
# RULES.md 準拠 定数定義（変更禁止）
# ============================================================
NAVY  = RGBColor(0x1A, 0x2B, 0x4A)
GREEN = RGBColor(0x00, 0x6B, 0x3F)
RED   = RGBColor(0xC4, 0x1E, 0x3A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY = RGBColor(0xE0, 0xE0, 0xE0)
BGRAY = RGBColor(0xF5, 0xF7, 0xFA)
GOLD  = RGBColor(0xB8, 0x86, 0x00)

SLIDE_W     = Cm(25.4)
SLIDE_H     = Cm(14.29)
MARGIN      = Cm(1.5)
GAP         = Cm(1.0)
TITLE_BAR_H = Cm(2.8)
ACCENT_H    = Pt(3)

# ============================================================
# 描画ヘルパー
# ============================================================

def rect(slide, l, t, w, h, fill=None, line=None, lw=Pt(1)):
    s = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        s.line.width = lw
    else:
        s.line.fill.background()
    return s

def tb(slide, l, t, w, h, text, size=Pt(15), bold=False,
       color=None, align=PP_ALIGN.LEFT):
    color = color or NAVY
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = 'Meiryo'
    return box

def mtb(slide, l, t, w, h, lines, default_size=Pt(14), default_color=None):
    """複数行テキストボックス。lines は str または dict のリスト"""
    default_color = default_color or NAVY
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            d = {'text': item, 'bold': False, 'color': default_color, 'size': default_size}
        else:
            d = item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = d.get('text', '')
        run.font.size = d.get('size', default_size)
        run.font.bold = d.get('bold', False)
        run.font.color.rgb = d.get('color', default_color)
        run.font.name = 'Meiryo'
    return box

def title_bar(slide, title, subtitle=None):
    """ネイビータイトルバー＋グリーンアクセントライン"""
    rect(slide, 0, 0, SLIDE_W, TITLE_BAR_H, fill=NAVY)
    rect(slide, 0, TITLE_BAR_H, SLIDE_W, ACCENT_H, fill=GREEN)
    tb(slide, MARGIN, Cm(0.2), SLIDE_W - MARGIN*2, Cm(1.55),
       title, size=Pt(26), bold=True, color=WHITE)
    if subtitle:
        tb(slide, MARGIN, Cm(1.7), SLIDE_W - MARGIN*2, Cm(0.9),
           subtitle, size=Pt(13), color=LGRAY)
    return TITLE_BAR_H + ACCENT_H

def card_2x2(slide, bh, cards):
    """2×2グリッドレイアウト"""
    cw = (SLIDE_W - MARGIN*2 - GAP) / 2
    ch = (SLIDE_H - bh - MARGIN - GAP) / 2
    pos = [
        (MARGIN,              bh + Cm(0.4)),
        (MARGIN + cw + GAP,   bh + Cm(0.4)),
        (MARGIN,              bh + Cm(0.4) + ch + GAP),
        (MARGIN + cw + GAP,   bh + Cm(0.4) + ch + GAP),
    ]
    for (cx, cy), card in zip(pos, cards):
        tc = card.get('title_color', NAVY)
        rect(slide, cx, cy, cw, ch, fill=BGRAY, line=LGRAY, lw=Pt(1))
        rect(slide, cx, cy, cw, Cm(0.65), fill=tc)
        tb(slide, cx+Cm(0.2), cy+Cm(0.05), cw-Cm(0.4), Cm(0.6),
           card['title'], size=Pt(12), bold=True, color=WHITE)
        mtb(slide, cx+Cm(0.2), cy+Cm(0.75), cw-Cm(0.4), ch-Cm(0.85), card['lines'])

def two_col(slide, bh, left_lines, right_lines, split=0.5):
    """左右2カラムレイアウト"""
    mid = SLIDE_W * split
    rect(slide, mid, bh+Cm(0.3), Pt(2), SLIDE_H-bh-MARGIN, fill=LGRAY)
    lw = mid - MARGIN - Cm(0.3)
    rx = mid + Cm(0.8)
    rw = SLIDE_W - rx - MARGIN
    ch = SLIDE_H - bh - MARGIN
    mtb(slide, MARGIN,    bh+Cm(0.4), lw, ch, left_lines)
    mtb(slide, rx,        bh+Cm(0.4), rw, ch, right_lines)

def big_number_col(slide, bh, numbers, right_lines, split=0.38):
    """大数字＋説明テキストレイアウト"""
    mid = SLIDE_W * split
    rect(slide, mid, bh+Cm(0.3), Pt(2), SLIDE_H-bh-MARGIN, fill=LGRAY)
    lw = mid - MARGIN - Cm(0.3)
    rx = mid + Cm(0.8)
    rw = SLIDE_W - rx - MARGIN

    nh = (SLIDE_H - bh - MARGIN) / max(len(numbers), 1)
    for i, (label, val, unit, col) in enumerate(numbers):
        y = bh + Cm(0.1) + nh * i
        if i > 0:
            rect(slide, MARGIN, y, lw, Pt(1), fill=LGRAY)
        tb(slide, MARGIN, y+Cm(0.15), lw, Cm(0.65), label, size=Pt(13))
        tb(slide, MARGIN, y+Cm(0.75), lw, Cm(1.9),  val,   size=Pt(46), bold=True, color=col)
        if unit:
            tb(slide, MARGIN, y+Cm(2.55), lw, Cm(0.5), unit, size=Pt(13), color=col)

    mtb(slide, rx, bh+Cm(0.3), rw, SLIDE_H-bh-MARGIN-Cm(0.3), right_lines)

def vertical_list(slide, bh, items):
    """アイコン付き縦リスト（最大4項目）"""
    ih = (SLIDE_H - bh - MARGIN) / max(len(items), 1)
    for i, item in enumerate(items):
        y = bh + ih * i
        if i > 0:
            rect(slide, MARGIN, y, SLIDE_W-MARGIN*2, Pt(1.5), fill=LGRAY)
        icon      = item.get('icon', '●')
        title     = item.get('title', '')
        title_col = item.get('title_color', NAVY)
        lines     = item.get('lines', [])
        tb(slide, MARGIN+Cm(0.1), y+Cm(0.1), Cm(1.2), ih-Cm(0.2),
           icon, size=Pt(18))
        tb(slide, MARGIN+Cm(1.4), y+Cm(0.1), SLIDE_W-MARGIN*2-Cm(1.4), Cm(0.8),
           title, size=Pt(15), bold=True, color=title_col)
        for j, line in enumerate(lines):
            tb(slide, MARGIN+Cm(1.7), y+Cm(0.85)+Cm(0.5)*j,
               SLIDE_W-MARGIN*2-Cm(1.9), Cm(0.55), line, size=Pt(13))

def layout_e(slide, title, subtitle, topics):
    """1枚目専用トピック整理"""
    rect(slide, 0, 0, SLIDE_W, Cm(3.3), fill=NAVY)
    rect(slide, 0, Cm(3.3), SLIDE_W, Pt(4), fill=GREEN)
    tb(slide, MARGIN, Cm(0.2), SLIDE_W-MARGIN*2, Cm(1.65),
       title, size=Pt(32), bold=True, color=WHITE)
    if subtitle:
        tb(slide, MARGIN, Cm(1.85), SLIDE_W-MARGIN*2, Cm(1.1),
           subtitle, size=Pt(14), color=LGRAY)
    top = Cm(3.5)
    ih = (SLIDE_H - top - MARGIN) / max(len(topics), 1)
    for i, topic in enumerate(topics):
        y = top + ih * i
        if i > 0:
            rect(slide, MARGIN, y, SLIDE_W-MARGIN*2, Pt(1), fill=LGRAY)
        tb(slide, MARGIN+Cm(0.3), y+Cm(0.1),
           SLIDE_W-MARGIN*2-Cm(0.6), ih-Cm(0.15), topic, size=Pt(16))

def quote_box(slide, x, y, w, h, text):
    """クレイマー引用ボックス（ネイビー背景・白テキスト）"""
    rect(slide, x, y, w, h, fill=NAVY)
    tb(slide, x+Cm(0.3), y+Cm(0.25), w-Cm(0.6), h-Cm(0.4),
       text, size=Pt(13), color=WHITE)

def disclaimer(slide):
    """最終スライド用免責注記"""
    ny = SLIDE_H - Cm(1.5)
    rect(slide, 0, ny, SLIDE_W, Cm(1.5), fill=BGRAY)
    tb(slide, MARGIN, ny+Cm(0.3), SLIDE_W-MARGIN*2, Cm(1.0),
       "※本スライドは情報提供のみを目的としており、投資の勧誘を意図するものではありません。",
       size=Pt(11))

# ============================================================
# plan.md パーサー
# ============================================================

def color_from_marker(text):
    """🟢🟡🔴 テキストからカラーを返す"""
    if '🟢' in text:
        return GREEN
    elif '🔴' in text:
        return RED
    elif '🟡' in text:
        return GOLD
    return NAVY

def line_to_dict(line):
    """1行テキストをスタイルdictに変換"""
    line = re.sub(r'^[-・\s]+', '', line).strip()
    bold = '**' in line
    text = line.replace('**', '').strip()
    # マイナス数値 → RED、プラス → GREEN
    if re.search(r'[−\-]\d', text):
        color = RED
    elif re.search(r'\+\d', text):
        color = GREEN
    else:
        color = NAVY
    return {'text': text, 'bold': bold, 'color': color, 'size': Pt(14)}

def parse_bullet_block(raw):
    """箇条書きブロック文字列 → dict リスト"""
    result = []
    for line in raw.strip().split('\n'):
        line = line.strip()
        if not line:
            result.append({'text': '', 'size': Pt(6), 'color': NAVY})
            continue
        # インデントされた行（説明補足）は少し小さく
        if line.startswith('  ') or line.startswith('\t'):
            d = line_to_dict(line)
            d['size'] = Pt(13)
            result.append(d)
        else:
            result.append(line_to_dict(line))
    return result

def parse_plan(plan_path):
    """plan.md を読み込んでスライドデータリストを返す"""
    with open(plan_path, encoding='utf-8') as f:
        content = f.read()

    # ## スライドN で分割（最初のメタブロックは除く）
    blocks = re.split(r'\n## スライド\d+[：:（(]?', content)[1:]

    slides = []
    for block in blocks:
        sd = {}

        # レイアウト
        m = re.search(r'\[LAYOUT-([A-E])\]', block)
        sd['layout'] = m.group(1) if m else 'A'

        # タイトル
        m = re.search(r'\*\*タイトル\*\*\s*[:：]\s*(.+)', block)
        sd['title'] = m.group(1).strip() if m else ''

        # サブタイトル
        m = re.search(r'\*\*サブ(?:タイトル)?\*\*\s*[:：]\s*(.+)', block)
        sd['subtitle'] = m.group(1).strip() if m else None

        # --- LAYOUT-E ---
        if sd['layout'] == 'E':
            topics = re.findall(r'^[①②③④⑤⑥⑦⑧⑨⑩].+', block, re.MULTILINE)
            sd['topics'] = [t.strip() for t in topics]

        # --- LAYOUT-C ---
        elif sd['layout'] == 'C':
            # 大数字ブロック
            m = re.search(r'\*\*大数字.*?\*\*[^\n]*\n(.*?)(?=\n\*\*説明|\Z)', block, re.DOTALL)
            numbers = []
            if m:
                for line in m.group(1).strip().split('\n'):
                    line = re.sub(r'^[-・\s]+', '', line.strip()).replace('**', '')
                    if not line:
                        continue
                    # "ラベル　値単位" の形式で分割
                    parts = re.split(r'[\s　]+', line, 1)
                    label = parts[0]
                    rest  = parts[1] if len(parts) > 1 else ''
                    # 値と単位を分離（例: −444ポイント → −444, ポイント）
                    m2 = re.match(r'([^\u3040-\u9fff\uff00-\uffefA-Za-z（）()]+)(.*)', rest)
                    val  = m2.group(1).strip() if m2 else rest
                    unit = m2.group(2).strip() if m2 else ''
                    col  = RED if ('−' in val or (val.startswith('-') and len(val) > 1)) else GREEN
                    numbers.append((label, val, unit, col))
            sd['numbers'] = numbers

            # 説明テキスト
            m = re.search(r'\*\*説明テキスト.*?\*\*[^\n]*\n(.*?)(?=\n---|\n##|\Z)', block, re.DOTALL)
            sd['right_lines'] = parse_bullet_block(m.group(1) if m else '')

            # 引用ボックス（「〜」形式）
            m = re.search(r'「([^」]{10,120})」', block)
            sd['quote'] = f'「{m.group(1)}」' if m else None

        # --- LAYOUT-A ---
        elif sd['layout'] == 'A':
            m = re.search(r'\*\*左カラム.*?\*\*[^\n]*\n(.*?)(?=\n\*\*右カラム|\Z)', block, re.DOTALL)
            sd['left_lines'] = parse_bullet_block(m.group(1) if m else '')

            m = re.search(r'\*\*右カラム.*?\*\*[^\n]*\n(.*?)(?=\n---|\n##|\*\*注記|\Z)', block, re.DOTALL)
            sd['right_lines'] = parse_bullet_block(m.group(1) if m else '')

            m = re.search(r'「([^」]{10,120})」', block)
            sd['quote'] = f'「{m.group(1)}」' if m else None

        # --- LAYOUT-B ---
        elif sd['layout'] == 'B':
            # カードヘッダーと本文を抽出
            pattern = r'\*\*(カード[①②③④][^\*]*)\*\*[^\n]*\n(.*?)(?=\n\*\*カード[①②③④]|\n---|\n##|\Z)'
            matches = re.findall(pattern, block, re.DOTALL)
            cards = []
            for header, body in matches:
                tc = color_from_marker(header)
                # タイトル整形：「カード①（〜）」→「〜」
                title = re.sub(r'カード[①②③④]', '', header)
                title = re.sub(r'^[（(]|[）)]$', '', title.strip()).strip()
                if not title:
                    title = header.strip()
                lines = parse_bullet_block(body)
                cards.append({'title': title, 'title_color': tc, 'lines': lines})
            sd['cards'] = cards

        # --- LAYOUT-D ---
        elif sd['layout'] == 'D':
            pattern = r'\*\*[①②③④]\s+(.+?)\*\*[^\n]*\n(.*?)(?=\n\*\*[①②③④]|\n---|\n##|\Z)'
            matches = re.findall(pattern, block, re.DOTALL)
            items = []
            for header, body in matches:
                tc   = color_from_marker(header)
                icon = '🟢' if tc == GREEN else ('🔴' if tc == RED else '🟡')
                lines = [re.sub(r'^[-・\s]+', '', l).strip()
                         for l in body.strip().split('\n') if l.strip()]
                items.append({'icon': icon, 'title': header.strip(),
                              'title_color': tc, 'lines': lines})
            sd['items'] = items

        # 最終スライド判定（注記テキスト or 最後のスライド）
        sd['has_disclaimer'] = '※本スライドは情報提供' in block

        slides.append(sd)

    # 最後のスライドには必ず免責注記
    if slides:
        slides[-1]['has_disclaimer'] = True

    return slides

# ============================================================
# スライド生成
# ============================================================

def build_slides(prs, slides_data):
    blank = prs.slide_layouts[6]

    for sd in slides_data:
        s      = prs.slides.add_slide(blank)
        layout = sd['layout']

        if layout == 'E':
            layout_e(s, sd['title'], sd.get('subtitle'), sd.get('topics', []))

        elif layout == 'C':
            bh = title_bar(s, sd['title'], sd.get('subtitle'))
            big_number_col(s, bh, sd.get('numbers', []), sd.get('right_lines', []))
            if sd.get('quote'):
                mid = SLIDE_W * 0.38
                rx  = mid + Cm(0.8)
                rw  = SLIDE_W - rx - MARGIN
                qy  = SLIDE_H - MARGIN - Cm(1.8)
                quote_box(s, rx, qy, rw, Cm(1.8), sd['quote'])

        elif layout == 'A':
            bh = title_bar(s, sd['title'], sd.get('subtitle'))
            two_col(s, bh, sd.get('left_lines', []), sd.get('right_lines', []))
            if sd.get('quote'):
                mid = SLIDE_W * 0.5
                rx  = mid + Cm(0.8)
                rw  = SLIDE_W - rx - MARGIN
                qy  = SLIDE_H - MARGIN - Cm(2.8)
                quote_box(s, rx, qy, rw, Cm(2.8), sd['quote'])

        elif layout == 'B':
            bh = title_bar(s, sd['title'], sd.get('subtitle'))
            card_2x2(s, bh, sd.get('cards', []))

        elif layout == 'D':
            bh = title_bar(s, sd['title'], sd.get('subtitle'))
            vertical_list(s, bh, sd.get('items', []))

        if sd.get('has_disclaimer'):
            disclaimer(s)

# ============================================================
# メイン
# ============================================================

def main():
    if len(sys.argv) < 2:
        print("使い方: python scripts/generate_pptx.py plans/plan_YYYYMMDD_タイトル.md")
        sys.exit(1)

    plan_path = sys.argv[1]
    if not os.path.exists(plan_path):
        print(f"エラー: ファイルが見つかりません → {plan_path}")
        sys.exit(1)

    # 出力ファイル名をplan名の日付から自動生成
    base       = os.path.basename(plan_path)
    m          = re.search(r'(\d{8})', base)
    date_str   = m.group(1) if m else 'YYYYMMDD'
    out_dir    = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(plan_path))), 'outputs')
    os.makedirs(out_dir, exist_ok=True)
    out_path   = os.path.join(out_dir, f'{date_str}_マッドマネー要約.pptx')

    print(f"plan読み込み中: {plan_path}")
    slides_data = parse_plan(plan_path)
    print(f"  → {len(slides_data)}枚のスライドを検出")

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    build_slides(prs, slides_data)
    prs.save(out_path)
    print(f"保存完了: {out_path}")

if __name__ == '__main__':
    main()
