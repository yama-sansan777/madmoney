from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# プレゼンテーションオブジェクトを作成
prs = Presentation()
# 16:9比率 幅25.4cm × 高さ14.29cm（RULES.md準拠）
prs.slide_width = Inches(10)  # 25.4cm ≈ 10 inches
prs.slide_height = Inches(5.63)  # 14.29cm ≈ 5.63 inches (16:9比率)

# カラーパレット（RULES.md準拠）
NAVY = RGBColor(26, 43, 74)  # ダークネイビー #1A2B4A
DEEP_GREEN = RGBColor(0, 107, 63)  # ディープグリーン #006B3F
ALERT_RED = RGBColor(196, 30, 58)  # アラートレッド #C41E3A
WHITE = RGBColor(255, 255, 255)  # ホワイト #FFFFFF
LIGHT_GRAY = RGBColor(224, 224, 224)  # ライトグレー #E0E0E0
BG_GRAY = RGBColor(245, 247, 250)  # 背景グレー #F5F7FA

def add_slide_with_header(prs, title_text):
    """ヘッダー付きスライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 上部ヘッダーバー（ネイビー）- 高さを16:9比率に調整
    header_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0), Inches(10), Inches(0.9)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = NAVY
    header_shape.line.fill.background()

    # タイトルテキスト（白文字）
    title_frame = header_shape.text_frame
    title_frame.text = title_text
    title_frame.word_wrap = True
    title_frame.margin_top = Inches(0.2)
    title_frame.margin_left = Inches(0.5)
    title_frame.margin_right = Inches(0.5)
    p = title_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    return slide

def add_content_slide(prs, title_text, content_items=None):
    """コンテンツスライドを追加（箇条書き）"""
    slide = add_slide_with_header(prs, title_text)

    # コンテンツ - 16:9比率に調整
    if content_items:
        content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(8.8), Inches(4.3))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for i, item in enumerate(content_items):
            if i > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[i]
            p.text = item
            p.font.size = Pt(15)
            p.space_after = Pt(8)
            p.level = 0
            # インデント調整
            if item.startswith("   "):
                p.level = 1

    return slide

# Slide 01: 本日の主要トピック整理
slide = add_content_slide(prs, "本日の主要トピック", [
    "① 市場の大逆転劇 ── トランプ投稿がすべてを変えた",
    "② 4市場が語る「和平票」── 株・原油・債券・金の動き",
    "③ スワーマー（SWRM）IPO ── 防衛ドローンSWの衝撃デビュー",
    "④ アルタ・ビューティー CEO インタビュー",
    "⑤ ジェイコブズ・ソリューションズ CEO インタビュー",
    "⑥ ライトニング・ラウンド：10銘柄の評価",
    "⑦ 本日の3つの教訓"
])

# Slide 02: 相場の大逆転：朝4時 vs 引け後
slide = add_slide_with_header(prs, "朝4時と引け後 ── 同じ日とは思えない景色")

# 左カラム背景 - はみ出さないように調整（画面高さ5.63 - ヘッダー0.9 = 4.73インチ利用可能）
left_bg = slide.shapes.add_shape(1, Inches(0.4), Inches(1.1), Inches(4.6), Inches(4.3))
left_bg.fill.solid()
left_bg.fill.fore_color.rgb = BG_GRAY
left_bg.line.color.rgb = LIGHT_GRAY

# 左カラムテキスト
left_box = slide.shapes.add_textbox(Inches(0.55), Inches(1.25), Inches(4.3), Inches(4))
left_frame = left_box.text_frame
left_frame.word_wrap = True
left_items = [
    "午前4時：最悪の展開",
    "",
    "🔴 金利：上昇傾向",
    "🔴 原油：100ドル超",
    "🔴 金：証拠金売りで急落",
    "🔴 株式先物：ナスダック▲1%",
    "🔴 S&P500：直近高値から▲7.6%",
    "",
    "クレイマー「空売り勢は天才に見えた」"
]
for i, item in enumerate(left_items):
    if i > 0:
        left_frame.add_paragraph()
    p = left_frame.paragraphs[i]
    p.text = item
    p.font.size = Pt(13)
    if i == 0:
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = ALERT_RED
    elif "クレイマー" in item:
        p.font.size = Pt(11)
        p.font.italic = True

# 右カラム背景
right_bg = slide.shapes.add_shape(1, Inches(5.1), Inches(1.1), Inches(4.5), Inches(4.3))
right_bg.fill.solid()
right_bg.fill.fore_color.rgb = BG_GRAY
right_bg.line.color.rgb = LIGHT_GRAY

# 右カラムテキスト
right_box = slide.shapes.add_textbox(Inches(5.25), Inches(1.25), Inches(4.2), Inches(4))
right_frame = right_box.text_frame
right_frame.word_wrap = True
right_items = [
    "引け：大逆転",
    "",
    "🟢 ダウ：+631ポイント",
    "🟢 S&P500：+1.15%",
    "🟢 ナスダック：+1.38%",
    "🟢 原油：84ドルまで急落",
    "",
    "午前7時5分、トランプ大統領のSNS投稿が転換点"
]
for i, item in enumerate(right_items):
    if i > 0:
        right_frame.add_paragraph()
    p = right_frame.paragraphs[i]
    p.text = item
    p.font.size = Pt(13)
    if i == 0:
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = DEEP_GREEN
    elif "トランプ" in item:
        p.font.size = Pt(11)
        p.font.italic = True

# Slide 03: 転換点：午前7時5分の投稿
slide = add_slide_with_header(prs, "すべてを変えた1つのSNS投稿")

# 大数字背景 - 16:9比率に調整
number_bg = slide.shapes.add_shape(1, Inches(3), Inches(1.3), Inches(4), Inches(1.2))
number_bg.fill.solid()
number_bg.fill.fore_color.rgb = BG_GRAY
number_bg.line.color.rgb = LIGHT_GRAY

# 大数字
big_number_box = slide.shapes.add_textbox(Inches(3), Inches(1.35), Inches(4), Inches(1.1))
big_number_frame = big_number_box.text_frame
big_number_frame.text = "7:05 AM"
big_number_frame.paragraphs[0].font.size = Pt(60)
big_number_frame.paragraphs[0].font.bold = True
big_number_frame.paragraphs[0].font.color.rgb = NAVY
big_number_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# 説明テキスト
content_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.7), Inches(8.8), Inches(1.8))
content_frame = content_box.text_frame
content_frame.word_wrap = True
content_items = [
    "• トランプ大統領がSNSにイランとの和平交渉進展を示唆する投稿",
    "• 原油価格が投稿直後に100ドル超 → 84ドルへ急落",
    "• ダウは一時+1,000ポイントまで急伸、終値+631ポイント",
    "• ナスダックは朝▲1% → 引け+1.38%へ劇的転換",
    "• 午前11時51分：イラン国会議長が「交渉などない」と否定投稿 → 市場は再反発"
]
for i, item in enumerate(content_items):
    if i > 0:
        content_frame.add_paragraph()
    p = content_frame.paragraphs[i]
    p.text = item
    p.font.size = Pt(12)
    p.space_after = Pt(6)

# クレイマーコメントボックス
comment_shape = slide.shapes.add_shape(1, Inches(0.6), Inches(4.7), Inches(8.8), Inches(0.7))
comment_shape.fill.solid()
comment_shape.fill.fore_color.rgb = NAVY
comment_shape.line.color.rgb = NAVY
comment_text = comment_shape.text_frame
comment_text.word_wrap = True
comment_text.margin_top = Inches(0.12)
comment_text.margin_left = Inches(0.3)
comment_text.text = "「今日の上昇は恐怖の上昇だ」── 実態を伴うかどうかは翌日の市場次第"
comment_text.paragraphs[0].font.size = Pt(13)
comment_text.paragraphs[0].font.color.rgb = WHITE
comment_text.paragraphs[0].font.italic = True

# Slide 04: 4市場が投票した「和平の可能性」
slide = add_slide_with_header(prs, "4市場の「有権者」はどう判断したか")

# 2x2グリッド - 16:9比率に調整
grid_data = [
    {
        "title": "原油市場",
        "emoji": "🟢",
        "subtitle": "最も早く信じた",
        "content": "100ドル超 → 84ドルへ急落\n停戦＝中東供給安定と判断",
        "pos": (0.5, 1.2, 4.6, 1.9)
    },
    {
        "title": "株式市場",
        "emoji": "🟢",
        "subtitle": "踏み上げで大幅高",
        "content": "空売り投資家が損失回避で買い戻し\nダウ+631、ナスダック+1.38%",
        "pos": (5.2, 1.2, 4.3, 1.9)
    },
    {
        "title": "債券市場",
        "emoji": "🟡",
        "subtitle": "反応は複雑・遅延",
        "content": "10年・20年金利は午前中ほぼ変化なし\n午前中頃にようやく低下→「利下げ可能性が戻ってきた」シグナル",
        "pos": (0.5, 3.3, 4.6, 1.9)
    },
    {
        "title": "金市場",
        "emoji": "🔴",
        "subtitle": "判断が難しい票",
        "content": "前夜から▲9%（4,098ドルまで下落）\n投稿の1時間半前に反転開始も終値プラス転換ならず",
        "pos": (5.2, 3.3, 4.3, 1.9)
    }
]

for card in grid_data:
    # カード背景
    card_bg = slide.shapes.add_shape(
        1, Inches(card["pos"][0]), Inches(card["pos"][1]),
        Inches(card["pos"][2]), Inches(card["pos"][3])
    )
    card_bg.fill.solid()
    card_bg.fill.fore_color.rgb = BG_GRAY
    card_bg.line.color.rgb = LIGHT_GRAY

    # カードテキスト
    box = slide.shapes.add_textbox(
        Inches(card["pos"][0] + 0.12), Inches(card["pos"][1] + 0.12),
        Inches(card["pos"][2] - 0.24), Inches(card["pos"][3] - 0.24)
    )
    text_frame = box.text_frame
    text_frame.word_wrap = True

    # タイトル行
    p = text_frame.paragraphs[0]
    p.text = f"{card['emoji']} {card['title']}"
    p.font.size = Pt(14)
    p.font.bold = True

    # サブタイトル
    text_frame.add_paragraph()
    p = text_frame.paragraphs[1]
    p.text = card["subtitle"]
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.space_after = Pt(4)

    # コンテンツ
    text_frame.add_paragraph()
    p = text_frame.paragraphs[2]
    p.text = card["content"]
    p.font.size = Pt(11)

# Slide 05: スワーマーIPO
slide = add_slide_with_header(prs, "スワーマー（SWRM） ── 公募5ドルから最大13倍、その後急落")

# 大数字背景 - 16:9比率に調整
number_bg = slide.shapes.add_shape(1, Inches(2.8), Inches(1.3), Inches(4.4), Inches(1.2))
number_bg.fill.solid()
number_bg.fill.fore_color.rgb = BG_GRAY
number_bg.line.color.rgb = LIGHT_GRAY

# 大数字
big_number_box = slide.shapes.add_textbox(Inches(2.8), Inches(1.35), Inches(4.4), Inches(0.9))
big_number_frame = big_number_box.text_frame
big_number_frame.text = "+520%"
big_number_frame.paragraphs[0].font.size = Pt(60)
big_number_frame.paragraphs[0].font.bold = True
big_number_frame.paragraphs[0].font.color.rgb = DEEP_GREEN
big_number_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ラベル
label_box = slide.shapes.add_textbox(Inches(2.8), Inches(2.3), Inches(4.4), Inches(0.3))
label_frame = label_box.text_frame
label_frame.text = "公募価格比（初日終値）"
label_frame.paragraphs[0].font.size = Pt(14)
label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
label_frame.paragraphs[0].font.color.rgb = NAVY

# 説明テキスト
content_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.8), Inches(8.8), Inches(1.5))
content_frame = content_box.text_frame
content_frame.word_wrap = True
content_items = [
    "• 公募価格：5ドル（先週月曜夜）",
    "• 初値：12.5ドル（公募比+150%）→ 初日終値：31ドル（+520%）",
    "• 翌水曜：65ドルの高値",
    "• 木曜：▲4.5%、金曜：▲30%超、放送当日：▲28%（26ドル台）",
    "• 🟢 それでも公募価格比で約5倍以上の水準を維持"
]
for i, item in enumerate(content_items):
    if i > 0:
        content_frame.add_paragraph()
    p = content_frame.paragraphs[i]
    p.text = item
    p.font.size = Pt(12)
    p.space_after = Pt(5)

# 会社概要ボックス
overview_bg = slide.shapes.add_shape(1, Inches(0.6), Inches(4.5), Inches(8.8), Inches(0.6))
overview_bg.fill.solid()
overview_bg.fill.fore_color.rgb = BG_GRAY
overview_bg.line.color.rgb = LIGHT_GRAY

overview_box = slide.shapes.add_textbox(Inches(0.75), Inches(4.6), Inches(8.5), Inches(0.4))
overview_frame = overview_box.text_frame
overview_frame.word_wrap = True
overview_frame.text = "会社概要：ドローン群の一括制御ソフトを開発する防衛テック企業。ウクライナ・ロシア戦争を契機に3年前設立。実戦で2年近く・10万回超の任務実績。"
overview_frame.paragraphs[0].font.size = Pt(11)
overview_frame.paragraphs[0].font.italic = True

# Slide 06: スワーマー：魅力 vs 財務の現実
slide = add_slide_with_header(prs, "魅力的な物語 vs 厳しい財務の現実")

# 左カラム背景 - はみ出さないように調整
left_bg = slide.shapes.add_shape(1, Inches(0.4), Inches(1.1), Inches(4.6), Inches(3.3))
left_bg.fill.solid()
left_bg.fill.fore_color.rgb = BG_GRAY
left_bg.line.color.rgb = LIGHT_GRAY

# 左カラムテキスト
left_box = slide.shapes.add_textbox(Inches(0.55), Inches(1.22), Inches(4.3), Inches(3.1))
left_frame = left_box.text_frame
left_frame.word_wrap = True
left_items = [
    "買われた理由",
    "",
    "🟢 ウクライナ戦争でドローン戦術の有効性が証明",
    "🟢 コスト非対称性「ホンダをランボルギーニにぶつける」戦術",
    "🟢 世界の軍が低コストドローン開発を競う時代",
    "🟢 エリック・プリンス氏が非常勤会長就任",
    "🟢 受注契約パイプライン：3,300万ドル"
]
for i, item in enumerate(left_items):
    if i > 0:
        left_frame.add_paragraph()
    p = left_frame.paragraphs[i]
    p.text = item
    if i == 0:
        p.font.bold = True
        p.font.size = Pt(15)
        p.font.color.rgb = DEEP_GREEN
    else:
        p.font.size = Pt(11)
        p.space_after = Pt(5)

# 右カラム背景
right_bg = slide.shapes.add_shape(1, Inches(5.1), Inches(1.1), Inches(4.5), Inches(3.3))
right_bg.fill.solid()
right_bg.fill.fore_color.rgb = BG_GRAY
right_bg.line.color.rgb = LIGHT_GRAY

# 右カラムテキスト
right_box = slide.shapes.add_textbox(Inches(5.22), Inches(1.22), Inches(4.3), Inches(3.1))
right_frame = right_box.text_frame
right_frame.word_wrap = True
right_items = [
    "慎重論の根拠",
    "",
    "🔴 売上高：31万ドル未満（昨年）",
    "🔴 損失：850万ドル超",
    "🔴 楽観試算でも売上高の16倍超の時価総額",
    "🔴 将来の株式希薄化リスク（SO未反映）",
    "🔴 引受証券会社が小規模ブティック1社のみ"
]
for i, item in enumerate(right_items):
    if i > 0:
        right_frame.add_paragraph()
    p = right_frame.paragraphs[i]
    p.text = item
    if i == 0:
        p.font.bold = True
        p.font.size = Pt(15)
        p.font.color.rgb = ALERT_RED
    else:
        p.font.size = Pt(11)
        p.space_after = Pt(5)

# クレイマーコメント
comment_shape = slide.shapes.add_shape(1, Inches(0.4), Inches(4.6), Inches(9.2), Inches(0.6))
comment_shape.fill.solid()
comment_shape.fill.fore_color.rgb = NAVY
comment_shape.line.color.rgb = NAVY
comment_text = comment_shape.text_frame
comment_text.word_wrap = True
comment_text.margin_top = Inches(0.1)
comment_text.margin_left = Inches(0.3)
comment_text.text = "「まだオーブンから出すには早すぎる」── 受注が実際の売上に結びつくか、株式数増加後の動きを見守る段階"
comment_text.paragraphs[0].font.size = Pt(12)
comment_text.paragraphs[0].font.color.rgb = WHITE

# Slide 07: アルタ・ビューティー
slide = add_slide_with_header(prs, "アルタ・ビューティー（ULTA） ── 株価急落後のCEOが語ること")

# 左カラム背景 - はみ出さないように調整
left_bg = slide.shapes.add_shape(1, Inches(0.4), Inches(1.1), Inches(4.6), Inches(2.7))
left_bg.fill.solid()
left_bg.fill.fore_color.rgb = BG_GRAY
left_bg.line.color.rgb = LIGHT_GRAY

# 左カラムテキスト
left_box = slide.shapes.add_textbox(Inches(0.55), Inches(1.22), Inches(4.3), Inches(2.5))
left_frame = left_box.text_frame
left_frame.word_wrap = True
left_items = [
    "足元の課題",
    "",
    "🔴 直近決算：コスト超過で利益が市場予想を下回る",
    "🔴 翌日株価：▲14%",
    "🔴 2月高値から▲28%"
]
for i, item in enumerate(left_items):
    if i > 0:
        left_frame.add_paragraph()
    p = left_frame.paragraphs[i]
    p.text = item
    if i == 0:
        p.font.bold = True
        p.font.size = Pt(15)
        p.font.color.rgb = ALERT_RED
    else:
        p.font.size = Pt(12)
        p.space_after = Pt(7)

# 右カラム背景
right_bg = slide.shapes.add_shape(1, Inches(5.1), Inches(1.1), Inches(4.5), Inches(2.7))
right_bg.fill.solid()
right_bg.fill.fore_color.rgb = BG_GRAY
right_bg.line.color.rgb = LIGHT_GRAY

# 右カラムテキスト
right_box = slide.shapes.add_textbox(Inches(5.22), Inches(1.22), Inches(4.3), Inches(2.5))
right_frame = right_box.text_frame
right_frame.word_wrap = True
right_items = [
    "成長戦略・CEO評価",
    "",
    "🟢 CEO就任約1年でS&P500を大幅アウトパフォーム",
    "🟢 SG&A増加率を売上成長率以内に抑えるとコミット",
    "🟢 TikTokショップへ参入（3月17日開始）",
    "🟢 GLP-1普及に伴う美容ニーズへの品揃え強化",
    "🟢 「消費の腰折れはまだ見えていない」"
]
for i, item in enumerate(right_items):
    if i > 0:
        right_frame.add_paragraph()
    p = right_frame.paragraphs[i]
    p.text = item
    if i == 0:
        p.font.bold = True
        p.font.size = Pt(15)
        p.font.color.rgb = DEEP_GREEN
    else:
        p.font.size = Pt(11)
        p.space_after = Pt(5)

# 補足テキスト背景
supplement_bg = slide.shapes.add_shape(1, Inches(0.4), Inches(4.0), Inches(9.2), Inches(0.45))
supplement_bg.fill.solid()
supplement_bg.fill.fore_color.rgb = BG_GRAY
supplement_bg.line.color.rgb = LIGHT_GRAY

supplement_box = slide.shapes.add_textbox(Inches(0.55), Inches(4.07), Inches(9), Inches(0.32))
supplement_frame = supplement_box.text_frame
supplement_frame.word_wrap = True
supplement_frame.text = "• CEO：キーシャ・スティールマン氏（2025年1月就任）  • クレイマー評価：現在の株価水準は「良い価格」── 引き続き支持"
supplement_frame.paragraphs[0].font.size = Pt(10)

# クレイマーコメント
comment_shape = slide.shapes.add_shape(1, Inches(0.4), Inches(4.6), Inches(9.2), Inches(0.6))
comment_shape.fill.solid()
comment_shape.fill.fore_color.rgb = NAVY
comment_shape.line.color.rgb = NAVY
comment_text = comment_shape.text_frame
comment_text.margin_top = Inches(0.1)
comment_text.margin_left = Inches(0.3)
comment_text.text = "「美容は景気後退に強いとは言えないが、耐性がある」"
comment_text.paragraphs[0].font.size = Pt(12)
comment_text.paragraphs[0].font.color.rgb = WHITE

# Slide 08: ジェイコブズ・ソリューションズ
slide = add_slide_with_header(prs, "ジェイコブズ（J） ── データセンター・リショアリング・生命科学")

# 2x2グリッド - はみ出さないように調整
grid_data = [
    {
        "title": "データセンター",
        "emoji": "🟢",
        "content": "過去1年で+62.2%成長\n受注パイプライン+500%\nNVIDIAのGTCに登壇",
        "pos": (0.4, 1.1, 4.7, 1.7)
    },
    {
        "title": "リショアリング",
        "emoji": "🟢",
        "content": "製造業の国内回帰需要が旺盛\n半導体工場・製薬施設の建設が中心\n大手製薬の米国内新設PJも受注",
        "pos": (5.2, 1.1, 4.4, 1.7)
    },
    {
        "title": "デジタルツイン技術",
        "emoji": "🟢",
        "content": "AIエージェントとの組み合わせ\n仮想空間でDC設計をシミュレーション\nGPU更新サイクルへの適応力が競争力",
        "pos": (0.4, 3.0, 4.7, 1.7)
    },
    {
        "title": "財務・評価",
        "emoji": "🟢",
        "content": "「ビート・アンド・レイズ」達成（2月決算）\n年初来▲2%は割安とクレイマー評価\n従業員5万人・常時7,000〜1万人の求人",
        "pos": (5.2, 3.0, 4.4, 1.7)
    }
]

for card in grid_data:
    # カード背景
    card_bg = slide.shapes.add_shape(
        1, Inches(card["pos"][0]), Inches(card["pos"][1]),
        Inches(card["pos"][2]), Inches(card["pos"][3])
    )
    card_bg.fill.solid()
    card_bg.fill.fore_color.rgb = BG_GRAY
    card_bg.line.color.rgb = LIGHT_GRAY

    # カードテキスト
    box = slide.shapes.add_textbox(
        Inches(card["pos"][0] + 0.12), Inches(card["pos"][1] + 0.12),
        Inches(card["pos"][2] - 0.24), Inches(card["pos"][3] - 0.24)
    )
    text_frame = box.text_frame
    text_frame.word_wrap = True

    # タイトル行
    p = text_frame.paragraphs[0]
    p.text = f"{card['emoji']} {card['title']}"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.space_after = Pt(4)

    # コンテンツ
    text_frame.add_paragraph()
    p = text_frame.paragraphs[1]
    p.text = card["content"]
    p.font.size = Pt(10)

# クレイマーコメント
comment_shape = slide.shapes.add_shape(1, Inches(0.4), Inches(4.9), Inches(9.2), Inches(0.6))
comment_shape.fill.solid()
comment_shape.fill.fore_color.rgb = NAVY
comment_shape.line.color.rgb = NAVY
comment_text = comment_shape.text_frame
comment_text.word_wrap = True
comment_text.margin_top = Inches(0.1)
comment_text.margin_left = Inches(0.3)
comment_text.text = "「いま語られるべきトレンドがすべてこの1社に集約されている」── 現在の株価は割安"
comment_text.paragraphs[0].font.size = Pt(12)
comment_text.paragraphs[0].font.color.rgb = WHITE

# Slide 09: ライトニング・ラウンド①
slide = add_slide_with_header(prs, "ライトニング・ラウンド ── クレイマーの即断評価（前半）")

# コンテンツ - 高さを調整してはみ出しを防ぐ
content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(8.8), Inches(4.3))
text_frame = content_box.text_frame
text_frame.word_wrap = True

lightning_items = [
    ("🟢 ボーイング（BA）── 「買い」", "title"),
    ("   受注残は見通せる限り消えない。戦時中の発注減への懸念は現実とかけ離れている", "detail"),
    ("🟢 RKT（ロケット・カンパニーズ）── 「買い」（14ドル近辺は割安）", "title"),
    ("   利下げ期待の代理指標。放送当日の午後にようやく利下げ可能性が意識され始めた", "detail"),
    ("🔴 サービスナウ（NOW）── 「まだ早い」", "title"),
    ("   PER26倍は安いが、ウォール街でSaaS全般が嫌われている局面。もう少し下落が続く可能性", "detail"),
    ("🟢 ニューモント・マイニング（NEM）── 「買い」（同業アグニコ・イーグルを好む）", "title"),
    ("   金価格は戦争開始以来▲16%。逆張りの好機", "detail"),
    ("🔴 サウンドハウンド（SOUN）── 「非推奨」", "title"),
    ("   継続的な赤字企業は推奨できないというクレイマーの一貫した立場", "detail")
]

for i, (item, item_type) in enumerate(lightning_items):
    if i > 0:
        text_frame.add_paragraph()
    p = text_frame.paragraphs[i]
    p.text = item
    if item_type == "title":
        p.font.size = Pt(13)
        p.font.bold = True
        p.space_after = Pt(2)
    else:  # detail
        p.font.size = Pt(11)
        p.space_after = Pt(6)

# Slide 10: ライトニング・ラウンド②
slide = add_slide_with_header(prs, "ライトニング・ラウンド ── クレイマーの即断評価（後半）")

# コンテンツ
content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(8.8), Inches(4.3))
text_frame = content_box.text_frame
text_frame.word_wrap = True

lightning_items2 = [
    ("🔴 フレッシュペット（FRPT）── 「様子見」", "title"),
    ("   利益率改善傾向も今回決算が市場期待に届かず。株価が大きく下落", "detail"),
    ("🔴 IEP（アイカーン・エンタープライジズ）── 「弱気（一貫）」", "title"),
    ("   空売り勢が正しかったという判断を維持", "detail"),
    ("🟢 ヴィアビ・ソリューションズ（VIAV）── 「前向き」（少し冷めてから買う）", "title"),
    ("   光ファイバーインフラ整備は最も活発な投資テーマ。株価が若干上がっているので冷めるのを待つ", "detail"),
    ("🔴 ホイールプール（WHR）── 「パス」", "title"),
    ("   高関税の恩恵を受けるはずが株価は▲25%超。下がり続ける株には買いを推奨できない", "detail"),
    ("", "spacer"),
    ("計10銘柄の内訳 → 🟢 買い・前向き：5銘柄 ／ 🔴 売り・非推奨・様子見：5銘柄", "summary")
]

for i, (item, item_type) in enumerate(lightning_items2):
    if i > 0:
        text_frame.add_paragraph()
    p = text_frame.paragraphs[i]
    p.text = item
    if item_type == "title":
        p.font.size = Pt(13)
        p.font.bold = True
        p.space_after = Pt(2)
    elif item_type == "detail":
        p.font.size = Pt(11)
        p.space_after = Pt(6)
    elif item_type == "spacer":
        p.font.size = Pt(6)
        p.space_after = Pt(0)
    else:  # summary
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.space_before = Pt(4)

# Slide 11: 本日の3つの教訓
slide = add_slide_with_header(prs, "今日の放送から投資家が学べること")

# コンテンツ
content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(8.8), Inches(4.3))
text_frame = content_box.text_frame
text_frame.word_wrap = True

lesson_items = [
    ("📊 株価だけ見ていると相場の本質を見誤る", "title"),
    ("   この日ダウは+631ポイント。しかし債券市場は午前中まで和平をほぼ織り込まず。", "detail"),
    ("   「株が上がった」事実だけでなく、債券・原油が何を言っているか合わせて確認する習慣が重要", "detail"),
    ("", "spacer"),
    ("💡 「良い話」と「良い株」は別物", "title"),
    ("   スワーマーの物語は本物。しかし売上高31万ドル未満で時価総額が売上高の16倍超。", "detail"),
    ("   テーマ株に熱くなる局面だからこそ、目論見書の財務数字を確認する習慣が大切", "detail"),
    ("", "spacer"),
    ("🛒 「悪いニュース」と「悪い消費」を区別する", "title"),
    ("   ニュースは暗い話題が絶えないが、ウォルマートは年初来+8%、ターゲットは+18%、", "detail"),
    ("   アルタのCEOは「値下がり商品への買い替えはまだ見えていない」と語った。", "detail"),
    ("   レジで何が売れているかは、ニュースより正直", "detail")
]

for i, (item, item_type) in enumerate(lesson_items):
    if i > 0:
        text_frame.add_paragraph()
    p = text_frame.paragraphs[i]
    p.text = item
    if item_type == "title":
        p.font.size = Pt(14)
        p.font.bold = True
        p.space_after = Pt(3)
    elif item_type == "detail":
        p.font.size = Pt(11)
        p.space_after = Pt(2)
    else:  # spacer
        p.font.size = Pt(8)
        p.space_after = Pt(0)

# Slide 12: 免責事項・チャンネル案内
slide = add_content_slide(prs, "ご視聴ありがとうございました", [
    "📋 本日の動画まとめ",
    "   テキスト版をメンバーシップ掲示板で公開中（最新動画の配信前に先行公開）",
    "",
    "🔔 チャンネル登録・いいねのお願い",
    "   動画が参考になった方はチャンネル登録・いいねで応援をお願いします",
    "",
    "🔒 メンバー限定コンテンツ",
    "   メンバー限定動画も配信中。気が向いたら覗いてみてください",
    "",
    "",
    "※本スライドは情報提供のみを目的としており、投資の勧誘を意図するものではありません。"
])

# 保存
import os

# outputsフォルダを作成（存在しない場合）
output_dir = r"c:\Users\TS2\Desktop\madmoney\madmoney\outputs"
os.makedirs(output_dir, exist_ok=True)

output_path = os.path.join(output_dir, "マッドマネー_20260324_fixed.pptx")
prs.save(output_path)
print(f"プレゼンテーションが作成されました: {output_path}")
